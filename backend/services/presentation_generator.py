import json
import logging
import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from config import get_settings
from services.vector_store import get_or_create_collection, ingest_documents
from models.schemas import PresentationConfig
import google.generativeai as genai
from fastapi.concurrency import run_in_threadpool

logger = logging.getLogger(__name__)


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert '#RRGGBB' string to pptx RGBColor."""
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _add_shadow(shape):
    """Add a subtle drop shadow to a shape (card effect)."""
    sp = shape._element
    spPr = sp.find(qn("a:spPr"))
    if spPr is None:
        spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        return
    effectLst = spPr.makeelement(qn("a:effectLst"), {})
    outerShdw = effectLst.makeelement(qn("a:outerShdw"), {
        "blurRad": "76200",    # 6pt blur
        "dist": "38100",       # 3pt distance
        "dir": "5400000",      # 270 degrees (below)
        "algn": "tl",
        "rotWithShape": "0",
    })
    srgbClr = outerShdw.makeelement(qn("a:srgbClr"), {"val": "000000"})
    alpha = srgbClr.makeelement(qn("a:alpha"), {"val": "25000"})  # 25% opacity
    srgbClr.append(alpha)
    outerShdw.append(srgbClr)
    effectLst.append(outerShdw)
    spPr.append(effectLst)


class PresentationGenerator:
    def __init__(self):
        self.settings = get_settings()
        if self.settings.gemini_api_key:
            genai.configure(api_key=self.settings.gemini_api_key)

    async def generate_presentation(
        self,
        session_id: str,
        topic: str = "Project Overview",
        num_slides: int = 10,
        config: PresentationConfig | None = None,
    ) -> tuple[str, list, PresentationConfig]:
        """
        Generates a PPTX file using a 2-stage AI pipeline (Flash -> Pro).
        Returns (file_path, slides_data, config_used).
        """
        cfg = config or PresentationConfig()

        # 1. Gather Context (RAG)
        context = await run_in_threadpool(self._get_global_context, session_id)
        if not context:
            raise ValueError("No documentation found. Please upload docs first.")

        # 2a. Agent 1: Content Extraction (Flash)
        key_points = await self._extract_key_points(context, topic)

        # 2b. Agent 2: Slide Structuring (Pro)
        slides_data = await self._structure_slides(key_points, topic, num_slides)

        # 3. Build premium PPTX
        pptx_path = await run_in_threadpool(self._build_pptx, slides_data, session_id, cfg)

        # 4. Auto-Ingest for voice bot
        await run_in_threadpool(self._ingest_presentation, pptx_path, slides_data, session_id)

        return pptx_path, slides_data, cfg

    # ── Context ──────────────────────────────────────────────────────────────

    def _get_global_context(self, session_id: str) -> str:
        """Fetch all global context documents for the session."""
        try:
            collection = get_or_create_collection(session_id)
            results = collection.get(where={"slide_id": "global"})
            docs = results.get("documents", [])
            full_context = "\n\n".join(docs)
            return full_context[:1000000]
        except Exception as e:
            logger.error(f"Error fetching context: {e}")
            return ""

    # ── Agent 1: Key Points ──────────────────────────────────────────────────

    async def _extract_key_points(self, context: str, topic: str) -> str:
        """Agent 1: Uses Flash model to extract key information."""
        system_prompt = """You are a Senior Research Analyst.
        Your task is to extract the most critical information from the provided documentation for a presentation.
        Focus on:
        - Problem Statement & Solution
        - Key Features & Benefits
        - Architecture & Technical Approach
        - Roadmap & Future Steps
        
        Output a structured summary of key points. Be comprehensive but concise.
        """
        try:
            model_name = getattr(self.settings, "gemini_flash_model", "gemini-2.0-flash")
            model = genai.GenerativeModel(model_name=model_name)
            response = await run_in_threadpool(
                model.generate_content,
                [system_prompt, f"Topic: {topic}\n\nContext:\n{context}"]
            )
            return response.text
        except Exception as e:
            logger.error(f"Agent 1 (Extraction) failed: {e}")
            raise ValueError(f"Content extraction failed: {e}")

    # ── Agent 2: Slide Structuring ───────────────────────────────────────────

    async def _structure_slides(self, key_points: str, topic: str, num_slides: int = 10) -> list:
        """Agent 2: Uses Pro model to structure the final slide deck JSON."""
        system_prompt = f"""You are an expert Presentation Architect.
        Create a professional slide deck JSON based on the provided key points.
        
        Output format: JSON Array of slide objects.
        Example:
        [
            {{
                "title": "Slide Title",
                "type": "title" | "content" | "section" | "agenda" | "architecture" | "closing",
                "subtitle": "Optional subtitle (for title/section slides)",
                "content": "Main content text. Use \\n for line breaks between bullet points.",
                "notes": "Detailed speaker notes for voiceover (conversational)."
            }}
        ]
        
        Guidelines:
        - Create exactly {num_slides} high-impact slides.
        - Storyline: Title -> Agenda -> Problem -> Solution -> Features -> Architecture -> Benefits -> Closing.
        - Each content slide should have 3-5 bullet points separated by \\n.
        - 'notes' must be scripted conversational text for a presenter to read.
        - title slide must have a "subtitle" field.
        - section slides act as dividers between major topics.
        """
        try:
            model_name = getattr(self.settings, "gemini_pro_model", "gemini-1.5-pro")
            model = genai.GenerativeModel(
                model_name=model_name,
                generation_config={"response_mime_type": "application/json"}
            )
            response = await run_in_threadpool(
                model.generate_content,
                [system_prompt, f"Topic: {topic}\n\nKey Points:\n{key_points}"]
            )
            return json.loads(response.text)
        except Exception as e:
            logger.error(f"Agent 2 (Structuring) failed: {e}")
            raise ValueError(f"Slide structuring failed: {e}")

    # ── Premium PPTX Builder ─────────────────────────────────────────────────

    def _build_pptx(self, slides_data: list, session_id: str, cfg: PresentationConfig) -> str:
        """Generate a premium consulting-grade PPTX with card-shadow layouts."""
        prs = Presentation()
        prs.slide_width = Inches(13.333)   # Widescreen 16:9
        prs.slide_height = Inches(7.5)

        for idx, slide_json in enumerate(slides_data):
            stype = slide_json.get("type", "content").lower()

            if stype == "title":
                self._add_title_slide(prs, slide_json, cfg)
            elif stype in ("section", "closing"):
                self._add_section_slide(prs, slide_json, cfg)
            elif stype == "agenda":
                self._add_agenda_slide(prs, slide_json, cfg)
            else:
                self._add_content_slide(prs, slide_json, cfg)

            # Speaker notes
            slide = prs.slides[-1]
            if slide_json.get("notes"):
                slide.notes_slide.notes_text_frame.text = slide_json["notes"]

            # Footer
            self._add_footer(slide, idx + 1, len(slides_data), cfg)

        # Save
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"ConsultDeck_{session_id}_{ts}.pptx"
        path = f"generated_presentations/{filename}"
        os.makedirs("generated_presentations", exist_ok=True)
        prs.save(path)
        logger.info(f"Generated premium PPTX: {path}")
        return path

    # ── Slide Builders ───────────────────────────────────────────────────────

    def _add_title_slide(self, prs, data, cfg):
        """Full-background title slide with accent bar and subtitle."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        bg_color = _hex_to_rgb(cfg.title_bg_color)
        accent = _hex_to_rgb(cfg.accent_color)
        fc = _hex_to_rgb(cfg.font_color)

        # Full background
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = bg_color

        # Top accent line
        line = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.08))
        line.fill.solid()
        line.fill.fore_color.rgb = accent
        line.line.fill.background()

        # Title
        tb = slide.shapes.add_textbox(Inches(1.5), Inches(2.0), Inches(10), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = data.get("title", "Presentation")
        p.font.size = Pt(cfg.font_size_title + 8)
        p.font.bold = True
        p.font.color.rgb = fc
        p.font.name = cfg.font_name

        # Accent bar
        bar = slide.shapes.add_shape(1, Inches(1.5), Inches(3.8), Inches(1.2), Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()

        # Subtitle
        sub_text = data.get("subtitle", data.get("content", ""))
        if sub_text:
            sb = slide.shapes.add_textbox(Inches(1.5), Inches(4.1), Inches(8), Inches(0.8))
            sf = sb.text_frame
            sf.word_wrap = True
            sp = sf.paragraphs[0]
            sp.text = sub_text
            sp.font.size = Pt(cfg.font_size_body + 2)
            sp.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)
            sp.font.name = cfg.font_name

        # Date
        db = slide.shapes.add_textbox(Inches(1.5), Inches(5.4), Inches(4), Inches(0.4))
        df = db.text_frame
        dp = df.paragraphs[0]
        dp.text = datetime.now().strftime("%B %Y")
        dp.font.size = Pt(cfg.font_size_caption + 2)
        dp.font.color.rgb = RGBColor(0x88, 0x99, 0xAA)
        dp.font.name = cfg.font_name

    def _add_content_slide(self, prs, data, cfg):
        """Content slide with colored header + white card with shadow."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        header_c = _hex_to_rgb(cfg.header_color)
        accent_c = _hex_to_rgb(cfg.accent_color)
        bg_c = _hex_to_rgb(cfg.background_color)
        fc = _hex_to_rgb(cfg.font_color)
        bfc = _hex_to_rgb(cfg.body_font_color)

        # Colored full background (like GenSpark)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = header_c

        # Title text in header area
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = data.get("title", "")
        p.font.size = Pt(cfg.font_size_heading + 4)
        p.font.bold = True
        p.font.color.rgb = fc
        p.font.name = cfg.font_name
        p.alignment = PP_ALIGN.CENTER

        # ── Card with shadow (GenSpark-style) ────────────────────────────
        card = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(1.0), Inches(1.5),
            Inches(11.333), Inches(5.2),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = bg_c
        card.line.fill.background()

        # Add drop shadow
        _add_shadow(card)

        # ── Bullet points inside card ────────────────────────────────────
        content = data.get("content", "")
        if content:
            lines = [l.strip() for l in content.split("\n") if l.strip()]

            y_offset = 1.8
            for i, line in enumerate(lines):
                clean = line.lstrip("•-*· 0123456789.)").strip()
                if not clean:
                    continue

                # Gold accent bullet
                bullet_box = slide.shapes.add_textbox(
                    Inches(1.5), Inches(y_offset), Inches(0.4), Inches(0.5)
                )
                bf = bullet_box.text_frame
                bp = bf.paragraphs[0]
                bp.text = "●"
                bp.font.size = Pt(12)
                bp.font.color.rgb = accent_c
                bp.font.name = cfg.font_name
                bp.font.bold = True

                # Bullet text
                text_box = slide.shapes.add_textbox(
                    Inches(2.0), Inches(y_offset), Inches(9.8), Inches(0.7)
                )
                ttf = text_box.text_frame
                ttf.word_wrap = True
                tp = ttf.paragraphs[0]
                tp.text = clean
                tp.font.size = Pt(cfg.font_size_body + 2)
                tp.font.color.rgb = bfc
                tp.font.name = cfg.font_name

                y_offset += 0.8

    def _add_section_slide(self, prs, data, cfg):
        """Section divider — full colored background with centered text."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        header_c = _hex_to_rgb(cfg.header_color)
        accent_c = _hex_to_rgb(cfg.accent_color)
        fc = _hex_to_rgb(cfg.font_color)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = header_c

        # Accent line
        line = slide.shapes.add_shape(
            1, Inches(5.5), Inches(2.8), Inches(2.3), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = accent_c
        line.line.fill.background()

        # Title
        tb = slide.shapes.add_textbox(Inches(1.5), Inches(3.1), Inches(10), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = data.get("title", "")
        p.font.size = Pt(cfg.font_size_title)
        p.font.bold = True
        p.font.color.rgb = fc
        p.font.name = cfg.font_name
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        sub = data.get("subtitle", data.get("content", ""))
        if sub:
            sb = slide.shapes.add_textbox(Inches(2.5), Inches(4.3), Inches(8), Inches(0.6))
            sf = sb.text_frame
            sf.word_wrap = True
            sp = sf.paragraphs[0]
            sp.text = sub
            sp.font.size = Pt(cfg.font_size_body)
            sp.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)
            sp.font.name = cfg.font_name
            sp.alignment = PP_ALIGN.CENTER

    def _add_agenda_slide(self, prs, data, cfg):
        """Agenda slide — colored bg + white cards for each item."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        header_c = _hex_to_rgb(cfg.header_color)
        accent_c = _hex_to_rgb(cfg.accent_color)
        bg_c = _hex_to_rgb(cfg.background_color)
        fc = _hex_to_rgb(cfg.font_color)
        bfc = _hex_to_rgb(cfg.body_font_color)

        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = header_c

        # Title
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = data.get("title", "Agenda")
        p.font.size = Pt(cfg.font_size_heading + 4)
        p.font.bold = True
        p.font.color.rgb = fc
        p.font.name = cfg.font_name
        p.alignment = PP_ALIGN.CENTER

        # Agenda items as individual cards
        content = data.get("content", "")
        items = [l.strip() for l in content.split("\n") if l.strip()]

        y_pos = 1.6
        for i, item in enumerate(items):
            clean = item.lstrip("0123456789.)-•*· ").strip()
            if not clean:
                continue

            # Item card
            card = slide.shapes.add_shape(
                1, Inches(1.5), Inches(y_pos), Inches(10.3), Inches(0.65)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = bg_c
            card.line.fill.background()
            _add_shadow(card)

            # Number badge
            num_box = slide.shapes.add_textbox(
                Inches(1.8), Inches(y_pos + 0.08), Inches(0.6), Inches(0.5)
            )
            nf = num_box.text_frame
            np = nf.paragraphs[0]
            np.text = f"{i + 1:02d}"
            np.font.size = Pt(20)
            np.font.bold = True
            np.font.color.rgb = accent_c
            np.font.name = cfg.font_name

            # Accent bar
            bar = slide.shapes.add_shape(
                1, Inches(2.5), Inches(y_pos + 0.25), Inches(0.4), Inches(0.04)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = accent_c
            bar.line.fill.background()

            # Item text
            item_box = slide.shapes.add_textbox(
                Inches(3.1), Inches(y_pos + 0.08), Inches(8.5), Inches(0.5)
            )
            itf = item_box.text_frame
            itf.word_wrap = True
            ip = itf.paragraphs[0]
            ip.text = clean
            ip.font.size = Pt(cfg.font_size_body + 2)
            ip.font.color.rgb = bfc
            ip.font.name = cfg.font_name

            y_pos += 0.85

    # ── Footer ───────────────────────────────────────────────────────────────

    def _add_footer(self, slide, num: int, total: int, cfg):
        """Brand footer + slide number + accent line."""
        prs_w = slide.part.package.presentation.slide_width
        accent_c = _hex_to_rgb(cfg.accent_color)

        # Bottom accent line
        line = slide.shapes.add_shape(1, Inches(0), Inches(7.3), prs_w, Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = accent_c
        line.line.fill.background()

        # Slide number
        nb = slide.shapes.add_textbox(Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.3))
        nf = nb.text_frame
        np = nf.paragraphs[0]
        np.text = f"{num} / {total}"
        np.font.size = Pt(9)
        np.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        np.font.name = cfg.font_name
        np.alignment = PP_ALIGN.RIGHT

        # Brand
        bb = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(2), Inches(0.3))
        bf = bb.text_frame
        bp = bf.paragraphs[0]
        bp.text = "ConsultDeck Studio"
        bp.font.size = Pt(8)
        bp.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
        bp.font.name = cfg.font_name

    # ── Ingest ───────────────────────────────────────────────────────────────

    def _ingest_presentation(self, file_path: str, slides_data: list, session_id: str):
        """Auto-ingest the generated slides for the voice bot."""
        for i, slide in enumerate(slides_data):
            doc = {
                "path": f"slide:gen_{i}:{slide.get('title','Slide')}",
                "content": f"Slide: {slide.get('title','')}\nType: {slide.get('type')}\nContent: {slide.get('content')}\nNotes: {slide.get('notes')}",
                "extension": ".slide"
            }
            ingest_documents(
                session_id=session_id,
                documents=[doc],
                slide_id=f"slide_gen_{i}"
            )
        logger.info(f"Auto-ingested {len(slides_data)} generated slides.")
