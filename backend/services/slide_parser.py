import fitz  # PyMuPDF
import io
import json
import logging
import asyncio
from typing import List, Dict, Any
from PIL import Image
from concurrent.futures import ThreadPoolExecutor

# For Vision fallback
import google.generativeai as genai
from config import get_settings

logger = logging.getLogger(__name__)

# Try importing python-pptx
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

from fastapi.concurrency import run_in_threadpool

async def parse_file(file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
    """Parse file into slide objects."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if ext == "pptx":
        return await _parse_pptx(file_bytes)
    elif ext == "pdf":
        return await _parse_pdf(file_bytes)
    else:
        raise ValueError(f"Unsupported file type: .{ext}")

async def _parse_pptx(file_bytes: bytes) -> List[Dict[str, Any]]:
    """Parse PowerPoint PPTX file."""
    if not HAS_PPTX:
        raise ImportError("python-pptx is required. pip install python-pptx")
    
    settings = get_settings()
    has_gemini = False
    if settings.gemini_api_key:
        genai.configure(api_key=settings.gemini_api_key)
        has_gemini = True
        
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    prs = await run_in_threadpool(Presentation, io.BytesIO(file_bytes))
    
    slides = []
    vision_tasks = []
    
    for idx, slide in enumerate(prs.slides):
        title = slide.shapes.title.text.strip() if slide.shapes.title else ""
        text_parts = []
        images = []
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text and text != title:
                        text_parts.append(text)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                images.append(shape)
                
        content = "\n".join(text_parts)
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            
        slides.append({
            "id": f"slide_{idx}",
            "slide_index": idx,
            "title": title or f"Slide {idx + 1}",
            "type": _detect_slide_type(idx, title, content, len(prs.slides)),
            "content": content,
            "notes": notes,
        })
        
        if len(content) < 50 and images and has_gemini:
            largest_image = max(images, key=lambda s: len(s.image.blob))
            img_blob = largest_image.image.blob
            vision_tasks.append((idx, img_blob))

    if vision_tasks:
        logger.info(f"[SlideParser] {len(vision_tasks)} slides need Vision enhancement. Running in parallel...")
        semaphore = asyncio.Semaphore(5) # Limit concurrency
        
        async def _vision_worker(idx, blob):
            async with semaphore:
                try:
                    img = Image.open(io.BytesIO(blob))
                    model = genai.GenerativeModel(getattr(settings, "gemini_flash_model", "gemini-2.5-flash"))
                    # generate_content_async is preferred here
                    response = await run_in_threadpool(model.generate_content, [
                        "Analyze this presentation slide image. Extract ALL visible text and summarize visual diagrams. Return a JSON object with keys: title (string), content (string: detailed points from slide), notes (string: speaker notes), type (string). Output ONLY JSON.",
                        img
                    ])
                    return idx, response.text
                except Exception as e:
                    logger.error(f"Vision failed for slide {idx}: {e}")
                    return idx, None

        results = await asyncio.gather(*[_vision_worker(idx, blob) for idx, blob in vision_tasks])
        
        for idx, text_resp in results:
            if not text_resp: continue
            try:
                if "```json" in text_resp:
                    text_resp = text_resp.split("```json")[1].split("```")[0]
                elif "```" in text_resp:
                    text_resp = text_resp.split("```")[1].split("```")[0]
                
                data = json.loads(text_resp.strip())
                slide = slides[idx]
                if not slide["title"] or slide["title"].startswith("Slide "):
                    slide["title"] = data.get("title", slide["title"])
                
                v_content = data.get("content", "")
                if isinstance(v_content, list): v_content = "\n".join([str(i) for i in v_content])
                if v_content: slide["content"] = str(v_content).strip()
                
                v_notes = data.get("notes", "")
                if isinstance(v_notes, list): v_notes = "\n".join([str(i) for i in v_notes])
                if v_notes: slide["notes"] = str(v_notes).strip()
                
            except Exception as e:
                logger.warning(f"Failed to parse Vision JSON for slide {idx}: {e}")

    return slides

async def _parse_pdf(file_bytes: bytes) -> List[Dict[str, Any]]:
    """Parse PDF file. Falls back to Vision API if text is missing."""
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    slides = []
    settings = get_settings()
    has_gemini = settings.gemini_api_key is not None
    
    vision_tasks = []
    
    for idx, page in enumerate(doc):
        text = page.get_text("text").strip()
        lines = [l.strip() for l in text.split("\n") if l.strip()]
        title = lines[0] if lines else f"Page {idx + 1}"
        content = "\n".join(lines[1:]) if len(lines) > 1 else text
        
        slides.append({
            "id": f"slide_{idx}",
            "slide_index": idx,
            "title": title[:100],
            "type": _detect_slide_type(idx, title, content, len(doc)),
            "content": content,
            "notes": "",
        })
        
        if len(text) < 50 and has_gemini:
            pix = page.get_pixmap(dpi=150)
            img_bytes = pix.tobytes("png")
            vision_tasks.append((idx, img_bytes))

    if vision_tasks:
        logger.info(f"[SlideParser] {len(vision_tasks)} PDF pages need Vision analysis.")
        semaphore = asyncio.Semaphore(5)
        
        async def _vision_pdf_worker(idx, img_data):
            async with semaphore:
                try:
                    img = Image.open(io.BytesIO(img_data))
                    model = genai.GenerativeModel(getattr(settings, "gemini_flash_model", "gemini-2.5-flash"))
                    response = await run_in_threadpool(model.generate_content, [
                        "Analyze this presentation slide. Extract ALL visible text and summarize visual diagrams. Return a JSON object with keys: title (string), content (string: detailed points from slide), notes (string: speaker notes), type (string). Output ONLY JSON.",
                        img
                    ])
                    return idx, response.text
                except Exception as e:
                    logger.error(f"PDF Vision failed for page {idx+1}: {e}")
                    return idx, None
        
        results = await asyncio.gather(*[_vision_pdf_worker(idx, data) for idx, data in vision_tasks])
        
        for idx, text_resp in results:
            if not text_resp: continue
            try:
                if "```json" in text_resp:
                    text_resp = text_resp.split("```json")[1].split("```")[0]
                elif "```" in text_resp:
                    text_resp = text_resp.split("```")[1].split("```")[0]
                
                data = json.loads(text_resp.strip())
                slide = slides[idx]
                slide["title"] = data.get("title", slide["title"])
                v_content = data.get("content", "")
                if isinstance(v_content, list): v_content = "\n".join([str(i) for i in v_content])
                if v_content: slide["content"] = str(v_content).strip()
                v_notes = data.get("notes", "")
                if isinstance(v_notes, list): v_notes = "\n".join([str(i) for i in v_notes])
                slide["notes"] = str(v_notes).strip()
                slide["type"] = data.get("type", slide["type"]).lower()
            except Exception as e:
                logger.warning(f"Failed to parse PDF Vision JSON for page {idx+1}: {e}")

    doc.close()
    return slides

def _detect_slide_type(index: int, title: str, content: str, total: int) -> str:
    title_lower = title.lower()
    if index == 0: return "title"
    if index == total - 1 and any(x in title_lower for x in ["thank", "q&a"]): return "closing"
    if len(content) < 50 and title: return "section"
    return "content"
